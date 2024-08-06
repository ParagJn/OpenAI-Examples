"""
Program Overview:
This Streamlit application is designed to summarize content from uploaded PDF files and provide the summary
in either a DOCX or on-screen display format. The application leverages the OpenAI API for text summarization.
Users can select from different summarization models and specify the number of pages they wish to summarize.
The summarized content can then be saved as a Word document or displayed directly on the Streamlit interface.

Key Features:
- Upload and process PDF files.
- Generate text summaries using OpenAI models.
- Save summaries as DOCX or show them on the screen in real-time.
- Streamlit-based user interface with customizable model and document type options.

Dependencies:
- os
- streamlit
- PyPDF2
- docx
- pptx
- io
- time

Author: parag.jn@gmail.com
Date: August 2024
"""

import os
import streamlit as st
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from openai_connector import OpenAIConnector
import io
import time

# Initialize OpenAI Connector
openai_connector = OpenAIConnector()

# Setting page config
st.set_page_config(
    page_title="Open Ai - Chatbot",
    page_icon="🧊",
    layout="wide",
    initial_sidebar_state="expanded",
    menu_items={'About': "# Content summarizer to word or powerpoint format!"}
)

def generate_summary_from_pdf(file, model, num_pages):
    """
    Generate a summary from a PDF document.

    This function reads the specified number of pages from an uploaded PDF document,
    then uses the OpenAI API to generate a concise summary highlighting key points.

    Args:
        file (UploadedFile): The uploaded PDF file.
        model (str): The name of the AI model to use for summarization.
        num_pages (int): The number of pages to summarize from the PDF.

    Returns:
        str: The generated summary as a string.

    Raises:
        StreamlitAPIException: If the number of pages is less than or equal to zero.
    """

    if num_pages > 0:
        prompt = "You are an expert summarizer. You need to summarize the document in a very concise manner highlighting key points in bulleted format"
        pdf = PdfReader(file)
        text = ""
        for page_num in range(min(num_pages, len(pdf.pages))):
            page = pdf.pages[page_num]
            text += page.extract_text()
        return openai_connector.summarize_text(text, model,prompt)
    else:
        st.warning("Number of pages needs to be greater than 0")
        st.stop()

def save_summary_as_docx(summary, filename):
    """
    Save a generated summary as a DOCX file.

    This function creates a DOCX document containing the provided summary with appropriate headings.

    Args:
        summary (str): The summary text to be saved.
        filename (str): The base file name to use for the saved DOCX file.

    Returns:
        str: The path to the saved DOCX file.
    """

    doc = Document()
    doc.add_heading('Document Summary', 0)
    doc.add_paragraph(summary)
    summary_path = f"{filename}-summary.docx"
    doc.save(summary_path)
    return summary_path

def save_summary_as_pptx(summary, filename):
    """
    Save a generated summary as a PPTX file.

    This function creates a PowerPoint presentation containing the provided summary on a single slide.

    Args:
        summary (str): The summary text to be saved.
        filename (str): The base file name to use for the saved PPTX file.

    Returns:
        str: The path to the saved PPTX file.
    """
    
    pres = Presentation()
    slide_layout = pres.slide_layouts[1]  # Title and Content layout
    slide = pres.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Document Summary"
    content.text = summary
    summary_path = f"{filename}-summary.pptx"
    pres.save(summary_path)
    return summary_path

# Streamlit UI
st.title("Content Summarizer")
st.write("---")

# Sidebar Inputs
model_choices = OpenAIConnector.model_choices
document_choices = ['Display on Screen', 'Generate DOCX']
model_type = st.sidebar.radio("Choose a model type:", model_choices)
doc_type = st.sidebar.radio("Choose document type:", document_choices)
num_pages_input = st.sidebar.text_input("Number of pages to summarize ", max_chars=5, value="1")

# File uploader
uploaded_file = st.file_uploader("Upload a PDF file", type="pdf")

def main():
    """
    Main function to handle UI interactions and generate document summaries.

    This function manages the Streamlit UI elements and coordinates the
    summarization and saving of PDF summaries based on user input.
    """

    def stream_data(data):
        """
        Stream data word-by-word with a slight delay (0.02 seconds).

        Args:
            data (str): The data string to be streamed.
        
        Yields:
            str: The next word in the data string to display.
        """
        for word in data.split(" "):
            yield word + " "
            time.sleep(0.02)
            
    if uploaded_file is not None:
        try:
            num_pages = int(num_pages_input) if num_pages_input else 1
            filename = os.path.splitext(uploaded_file.name)[0]
            
            with st.spinner("Generating summary..."):
                summary = generate_summary_from_pdf(uploaded_file, model_type, num_pages)
                summary = summary.replace('$','USD')
                
                if doc_type == 'Generate DOCX':
                    summary_file = save_summary_as_docx(summary, filename)
                    st.success(f"Summary generated and saved as {summary_file}")
                elif doc_type == "Display on Screen":
                    st.write_stream((stream_data(summary)))

        except Exception as e:
            st.error(f"Error: {str(e)}")

if __name__ == "__main__":
    main()